"""Document parsing by file extension."""

from __future__ import annotations

import hashlib
import io
import mimetypes
from pathlib import Path

import fitz  # pymupdf
from docx import Document
from pptx import Presentation

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".txt",
    ".md",
    ".png",
    ".jpg",
    ".jpeg",
}

MIN_PDF_TEXT_CHARS = 100


def _tesseract_available() -> bool:
    try:
        import pytesseract

        pytesseract.get_tesseract_version()
        return True
    except Exception:
        return False


def _ocr_image_bytes(image_bytes: bytes) -> str:
    import pytesseract
    from PIL import Image

    image = Image.open(io.BytesIO(image_bytes))
    return pytesseract.image_to_string(image).strip()


def _ocr_pdf_pages(doc: fitz.Document) -> str:
    parts: list[str] = []
    for page in doc:
        pix = page.get_pixmap(dpi=200)
        text = _ocr_image_bytes(pix.tobytes("png"))
        if text:
            parts.append(text)
    return "\n\n".join(parts).strip()


def _parse_pdf(data: bytes) -> tuple[str, int, bool, str]:
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        pages = doc.page_count
        parts = [page.get_text("text").strip() for page in doc]
        text = "\n\n".join(part for part in parts if part).strip()
        ocr_used = False

        if len(text) < MIN_PDF_TEXT_CHARS:
            if _tesseract_available():
                ocr_text = _ocr_pdf_pages(doc)
                if ocr_text:
                    text = ocr_text
                    ocr_used = True
            elif not text:
                raise RuntimeError(
                    "PDF has little extractable text and Tesseract OCR is not installed"
                )

        return text, pages, ocr_used, "pdf"
    finally:
        doc.close()


def _parse_docx(data: bytes) -> tuple[str, int, bool, str]:
    doc = Document(io.BytesIO(data))
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    text = "\n\n".join(parts).strip()
    return text, 1, False, "docx"


def _parse_pptx(data: bytes) -> tuple[str, int, bool, str]:
    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"[Notes] {notes}")
        if slide_parts:
            parts.append(f"## Slide {idx}\n\n" + "\n\n".join(slide_parts))
    text = "\n\n".join(parts).strip()
    return text, len(prs.slides), False, "pptx"


def _parse_text(data: bytes) -> tuple[str, int, bool, str]:
    text = data.decode("utf-8").strip()
    return text, 1, False, "text"


def _parse_image(data: bytes) -> tuple[str, int, bool, str]:
    if not _tesseract_available():
        raise RuntimeError(
            "Image OCR requires Tesseract. Install tesseract-ocr and ensure it is on PATH."
        )
    text = _ocr_image_bytes(data)
    return text, 1, True, "image"


def _to_markdown(text: str, parser: str) -> str:
    if not text:
        return ""
    if parser == "pptx":
        return text
    if parser == "text" and text.lstrip().startswith("#"):
        return text
    return text


def _guess_mime(filename: str) -> str:
    mime, _ = mimetypes.guess_type(filename)
    return mime or "application/octet-stream"


def parse_document(data: bytes, filename: str) -> dict:
    ext = Path(filename or "").suffix.lower()
    sha256 = hashlib.sha256(data).hexdigest()
    mime = _guess_mime(filename)

    base_meta = {
        "filename": filename or "upload",
        "mime": mime,
        "sha256": sha256,
        "pages": 0,
        "ocrUsed": False,
        "parser": "unknown",
    }

    if ext not in SUPPORTED_EXTENSIONS:
        return {
            "ok": False,
            "text": "",
            "markdown": "",
            "meta": base_meta,
            "error": (
                f"Unsupported file type '{ext or '(none)'}'. "
                f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
            ),
        }

    try:
        if ext == ".pdf":
            text, pages, ocr_used, parser = _parse_pdf(data)
        elif ext == ".docx":
            text, pages, ocr_used, parser = _parse_docx(data)
        elif ext == ".pptx":
            text, pages, ocr_used, parser = _parse_pptx(data)
        elif ext in {".txt", ".md"}:
            text, pages, ocr_used, parser = _parse_text(data)
        elif ext in {".png", ".jpg", ".jpeg"}:
            text, pages, ocr_used, parser = _parse_image(data)
        else:
            raise RuntimeError(f"No parser registered for {ext}")

        markdown = _to_markdown(text, parser)
        return {
            "ok": True,
            "text": text,
            "markdown": markdown,
            "meta": {
                **base_meta,
                "pages": pages,
                "ocrUsed": ocr_used,
                "parser": parser,
            },
        }
    except UnicodeDecodeError:
        return {
            "ok": False,
            "text": "",
            "markdown": "",
            "meta": {**base_meta, "parser": "text"},
            "error": "File is not valid UTF-8 text",
        }
    except Exception as exc:
        return {
            "ok": False,
            "text": "",
            "markdown": "",
            "meta": base_meta,
            "error": str(exc),
        }
